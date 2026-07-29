"""Render stored studio artifacts to downloadable office files (docx / pptx / pdf).

The frontend ``ExportButtons`` posts to ``/artifacts/{id}/export`` and then lets
the browser download the streamed bytes from the GET route. Formats per kind
mirror the UI's ``KIND_FORMATS``:

    report / test / card  → pdf, docx
    presentation          → pdf, pptx

Cyrillic PDFs need a Unicode TrueType font (the fpdf2 core fonts are latin-1).
We embed DejaVuSans (installed via ``fonts-dejavu-core`` in the Docker image);
on a dev host we fall back to any Unicode TTF we can find so the unit tests run.
The docx/pptx paths carry text as UTF-8 inside OOXML, so they need no font work.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

__all__ = ["ExportError", "export_artifact", "supported_formats"]


class ExportError(Exception):
    """Unsupported (kind, format) pair, or no PDF font available on this host."""


# --- format registry ------------------------------------------------------

MEDIA_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_FORMATS_BY_KIND: dict[str, set[str]] = {
    "report": {"pdf", "docx"},
    "test": {"pdf", "docx"},
    "card": {"pdf", "docx"},
    "presentation": {"pdf", "pptx"},
}


def supported_formats(kind: str) -> set[str]:
    """The export formats allowed for ``kind`` (empty for podcast/infographic)."""
    return set(_FORMATS_BY_KIND.get(kind, set()))


def export_artifact(kind: str, content: dict[str, Any], fmt: str) -> tuple[bytes, str]:
    """Render an artifact's ``content`` to ``fmt`` bytes.

    Returns ``(data, media_type)``. Raises :class:`ExportError` for an
    unsupported ``(kind, fmt)`` pair or a missing PDF font.
    """
    if fmt not in supported_formats(kind):
        raise ExportError(f"Cannot export kind={kind!r} as {fmt!r}")
    content = content or {}
    if fmt == "docx":
        return _build_docx(kind, content), MEDIA_TYPES["docx"]
    if fmt == "pptx":
        return _build_pptx(content), MEDIA_TYPES["pptx"]
    if fmt == "pdf":
        return _build_pdf(kind, content), MEDIA_TYPES["pdf"]
    raise ExportError(f"Unknown format {fmt!r}")  # pragma: no cover - guarded above


# --- DOCX (python-docx) ---------------------------------------------------


def _build_docx(kind: str, content: dict[str, Any]) -> bytes:
    from docx import Document

    doc = Document()
    if kind == "report":
        doc.add_heading("Отчёт", level=0)
        # v2 (reference structure): Резюме / Ключевой вывод / Риски / Рекомендации / Источники
        if content.get("summary"):
            doc.add_heading("Резюме", level=1)
            doc.add_paragraph(str(content["summary"]))
        if content.get("key_insight"):
            doc.add_heading("Ключевой вывод", level=1)
            doc.add_paragraph(str(content["key_insight"]))
        for heading, items in (("Риски", content.get("risks")), ("Рекомендации", content.get("recommendations"))):
            if items:
                doc.add_heading(heading, level=1)
                for i, item in enumerate(items, 1):
                    doc.add_paragraph(f"{i}. {item}")
        sources = content.get("sources") or []
        if sources:
            doc.add_heading("Источники", level=1)
            for i, src in enumerate(sources, 1):
                doc.add_paragraph(f"[{i}] {src}")
        # v1 (legacy stored artifacts)
        if content.get("executive_summary"):
            doc.add_heading("Резюме", level=1)
            doc.add_paragraph(str(content["executive_summary"]))
        metrics = content.get("metrics") or []
        if metrics:
            doc.add_heading("Ключевые метрики", level=1)
            for m in metrics:
                p = doc.add_paragraph(style="List Bullet")
                p.add_run(f"{m.get('label', '')}: ").bold = True
                p.add_run(str(m.get("value", "")))
        for s in content.get("sections") or []:
            doc.add_heading(str(s.get("heading", "")), level=1)
            doc.add_paragraph(str(s.get("body", "")))
    elif kind == "test":
        doc.add_heading("Тест", level=0)
        for i, q in enumerate(content.get("questions") or [], 1):
            doc.add_heading(f"{i}. {q.get('prompt', '')}", level=2)
            correct = q.get("correct")
            for j, opt in enumerate(q.get("options") or []):
                mark = "  ✓" if j == correct else ""
                doc.add_paragraph(f"{chr(65 + j)}. {opt}{mark}", style="List Bullet")
            if q.get("explanation"):
                p = doc.add_paragraph()
                p.add_run("Пояснение: ").bold = True
                p.add_run(str(q["explanation"]))
    elif kind == "card":
        doc.add_heading("Карточки с выводами", level=0)
        # v2 (reference deck): ИНСАЙТ / КОНТЕКСТ / ШАГи / РИСК / ДЕЙСТВИЕ
        for heading, text in (
            ("Инсайт", content.get("insight")),
            ("Контекст", content.get("context")),
        ):
            if text:
                doc.add_heading(heading, level=1)
                doc.add_paragraph(str(text))
        steps = [s for s in content.get("steps") or [] if isinstance(s, dict)]
        if steps:
            doc.add_heading("Шаги", level=1)
            for i, s in enumerate(steps, 1):
                p = doc.add_paragraph()
                p.add_run(f"Шаг {i}. {s.get('title', '')} — ").bold = True
                p.add_run(str(s.get("text", "")))
        for heading, text in (("Риск", content.get("risk")), ("Действие", content.get("action"))):
            if text:
                doc.add_heading(heading, level=1)
                doc.add_paragraph(str(text))
        # v1 (legacy stored artifacts)
        if content.get("title"):
            doc.add_heading(str(content["title"]), level=1)
        if content.get("summary"):
            doc.add_paragraph(str(content["summary"]))
        for kp in content.get("key_points") or []:
            p = doc.add_paragraph(style="List Bullet")
            p.add_run(f"{kp.get('label', '')}: ").bold = True
            p.add_run(str(kp.get("text", "")))
        recs = content.get("recommendations") or []
        if recs:
            doc.add_heading("Рекомендации", level=1)
            for r in recs:
                doc.add_paragraph(str(r), style="List Bullet")
        tags = content.get("tags") or []
        if tags:
            doc.add_paragraph("Теги: " + ", ".join(str(t) for t in tags))
    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# --- PPTX (python-pptx) ---------------------------------------------------


def _build_pptx(content: dict[str, Any]) -> bytes:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    title_layout = prs.slide_layouts[0]  # Title Slide
    bullet_layout = prs.slide_layouts[1]  # Title and Content

    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = str(content.get("title") or "Презентация")

    for idx, slide in enumerate(content.get("slides") or [], 1):
        s = prs.slides.add_slide(bullet_layout)
        heading = slide.get("heading") or slide.get("title") or f"Слайд {idx}"
        s.shapes.title.text = str(heading)
        body = s.placeholders[1].text_frame
        body.clear()
        bullets = slide.get("bullets") or []
        for k, b in enumerate(bullets):
            para = body.paragraphs[0] if k == 0 else body.add_paragraph()
            para.text = str(b)
            para.font.size = Pt(18)
        if slide.get("notes"):
            s.notes_slide.notes_text_frame.text = str(slide["notes"])

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- PDF (fpdf2 + embedded Unicode font) ----------------------------------

_REGULAR_FONT_CANDIDATES = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",  # Debian slim (our image)
    "/usr/share/fonts/dejavu/DejaVuSans.ttf",
    "/Library/Fonts/Arial Unicode.ttf",  # macOS
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",  # macOS dev host
)
_BOLD_FONT_CANDIDATES = (
    "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
    "/usr/share/fonts/dejavu/DejaVuSans-Bold.ttf",
)


def _first_existing(paths: tuple[str, ...]) -> str | None:
    return next((p for p in paths if Path(p).exists()), None)


def _new_pdf() -> Any:
    from fpdf import FPDF

    regular = _first_existing(_REGULAR_FONT_CANDIDATES)
    if regular is None:
        raise ExportError(
            "No Unicode TTF font available for PDF export "
            "(install fonts-dejavu-core in the image)."
        )
    bold = _first_existing(_BOLD_FONT_CANDIDATES) or regular
    pdf = FPDF()
    pdf.add_font("Body", "", regular)
    pdf.add_font("Body", "B", bold)
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Body", "", 12)
    return pdf


def _heading(pdf: Any, text: str, size: int = 15) -> None:
    # new_x=LMARGIN resets the cursor to the left margin (fpdf2 defaults to
    # new_x=RIGHT, which leaves width≈0 for the next full-width multi_cell).
    pdf.ln(2)
    pdf.set_font("Body", "B", size)
    pdf.multi_cell(0, 8, text, new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Body", "", 12)


def _para(pdf: Any, text: str) -> None:
    pdf.set_font("Body", "", 12)
    pdf.multi_cell(0, 6, text, new_x="LMARGIN", new_y="NEXT")
    pdf.ln(1)


def _bullet(pdf: Any, text: str) -> None:
    pdf.set_font("Body", "", 12)
    pdf.multi_cell(0, 6, f"•  {text}", new_x="LMARGIN", new_y="NEXT")


def _build_pdf(kind: str, content: dict[str, Any]) -> bytes:
    pdf = _new_pdf()
    if kind == "report":
        _heading(pdf, "Отчёт", size=20)
        # v2 (reference structure)
        if content.get("summary"):
            _heading(pdf, "Резюме")
            _para(pdf, str(content["summary"]))
        if content.get("key_insight"):
            _heading(pdf, "Ключевой вывод")
            _para(pdf, str(content["key_insight"]))
        for heading, items in (("Риски", content.get("risks")), ("Рекомендации", content.get("recommendations"))):
            if items:
                _heading(pdf, heading)
                for i, item in enumerate(items, 1):
                    _para(pdf, f"{i}. {item}")
        sources = content.get("sources") or []
        if sources:
            _heading(pdf, "Источники")
            for i, src in enumerate(sources, 1):
                _para(pdf, f"[{i}] {src}")
        # v1 (legacy stored artifacts)
        if content.get("executive_summary"):
            _heading(pdf, "Резюме")
            _para(pdf, str(content["executive_summary"]))
        metrics = content.get("metrics") or []
        if metrics:
            _heading(pdf, "Ключевые метрики")
            for m in metrics:
                _bullet(pdf, f"{m.get('label', '')}: {m.get('value', '')}")
        for s in content.get("sections") or []:
            _heading(pdf, str(s.get("heading", "")))
            _para(pdf, str(s.get("body", "")))
    elif kind == "test":
        _heading(pdf, "Тест", size=20)
        for i, q in enumerate(content.get("questions") or [], 1):
            _heading(pdf, f"{i}. {q.get('prompt', '')}")
            correct = q.get("correct")
            for j, opt in enumerate(q.get("options") or []):
                mark = "  ✓" if j == correct else ""
                _bullet(pdf, f"{chr(65 + j)}. {opt}{mark}")
            if q.get("explanation"):
                _para(pdf, f"Пояснение: {q['explanation']}")
    elif kind == "card":
        _heading(pdf, "Карточки с выводами", size=20)
        # v2 (reference deck)
        for heading, text in (
            ("Инсайт", content.get("insight")),
            ("Контекст", content.get("context")),
        ):
            if text:
                _heading(pdf, heading)
                _para(pdf, str(text))
        steps = [s for s in content.get("steps") or [] if isinstance(s, dict)]
        if steps:
            _heading(pdf, "Шаги")
            for i, s in enumerate(steps, 1):
                _para(pdf, f"Шаг {i}. {s.get('title', '')} — {s.get('text', '')}")
        for heading, text in (("Риск", content.get("risk")), ("Действие", content.get("action"))):
            if text:
                _heading(pdf, heading)
                _para(pdf, str(text))
        # v1 (legacy stored artifacts)
        if content.get("title"):
            _heading(pdf, str(content["title"]))
        if content.get("summary"):
            _para(pdf, str(content["summary"]))
        for kp in content.get("key_points") or []:
            _bullet(pdf, f"{kp.get('label', '')}: {kp.get('text', '')}")
        recs = content.get("recommendations") or []
        if recs:
            _heading(pdf, "Рекомендации")
            for r in recs:
                _bullet(pdf, str(r))
        tags = content.get("tags") or []
        if tags:
            _para(pdf, "Теги: " + ", ".join(str(t) for t in tags))
    elif kind == "presentation":
        _heading(pdf, str(content.get("title") or "Презентация"), size=20)
        for i, slide in enumerate(content.get("slides") or [], 1):
            heading = slide.get("heading") or slide.get("title") or f"Слайд {i}"
            _heading(pdf, f"{i}. {heading}")
            for b in slide.get("bullets") or []:
                _bullet(pdf, str(b))
            if slide.get("notes"):
                _para(pdf, str(slide["notes"]))
    return bytes(pdf.output())
